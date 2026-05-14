"""
EraseGraph — CMPE 273 Final Presentation
Completely redesigned: vivid color palette, 16 slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT = os.path.dirname(os.path.abspath(__file__))

# ── Color Palette (dark slate + vivid Tailwind-400 accents) ──────────────────
BG     = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
CARD   = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
CARD2  = RGBColor(0x26, 0x35, 0x4C)   # slate-750
WHITE  = RGBColor(0xF8, 0xFA, 0xFC)   # slate-50
SEC    = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
MUTED  = RGBColor(0x64, 0x74, 0x8B)   # slate-500
BORDER = RGBColor(0x33, 0x41, 0x55)   # slate-700

# Vivid 400-level accents — pop beautifully on dark backgrounds
ORANGE = RGBColor(0xFB, 0x92, 0x3C)   # orange-400
BLUE   = RGBColor(0x60, 0xA5, 0xFA)   # blue-400
GREEN  = RGBColor(0x34, 0xD3, 0x99)   # emerald-400
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)   # violet-400
AMBER  = RGBColor(0xFB, 0xBF, 0x24)   # amber-400
RED    = RGBColor(0xF8, 0x71, 0x71)   # red-400
TEAL   = RGBColor(0x22, 0xD3, 0xEE)   # cyan-400
PINK   = RGBColor(0xF4, 0x72, 0xB6)   # pink-400

# ── Canvas ────────────────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def rect(slide, x, y, w, h, fill, border=None, bw=0.8, radius=True):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    if radius:
        sp = shp.element
        prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 14000')
    return shp


def txt(slide, text, x, y, w, h, size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return tb


def image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, x, y, w, h)


def img_path(name):
    return os.path.join(OUT, name)


def accent_bar(slide, color=ORANGE):
    rect(slide, 0, 0, W, Inches(0.07), color, radius=False)


def slide_title(slide, text, color=WHITE, y=Inches(0.14), size=38):
    txt(slide, text, Inches(0.55), y, Inches(12.2), Inches(0.85),
        size=size, bold=True, color=color)


def slide_sub(slide, text, y=Inches(0.97), color=SEC):
    txt(slide, text, Inches(0.55), y, Inches(12.2), Inches(0.42),
        size=12.5, color=color)


def pill(slide, text, x, y, color, size=9.5):
    w = Inches(len(text) * 0.095 + 0.42)
    rect(slide, x, y, w, Inches(0.33), CARD2, color, 0.8)
    txt(slide, text, x + Inches(0.1), y + Inches(0.04), w - Inches(0.2),
        Inches(0.26), size=size, bold=True, color=color, align=PP_ALIGN.CENTER)
    return w


def card(slide, x, y, w, h, accent, title, bullets):
    rect(slide, x, y, w, h, CARD, accent, 1.0)
    txt(slide, title, x + Inches(0.2), y + Inches(0.17), w - Inches(0.4),
        Inches(0.34), size=12, bold=True, color=accent)
    cy = y + Inches(0.60)
    for b in bullets:
        txt(slide, '  ▸  ' + b, x + Inches(0.2), cy, w - Inches(0.4),
            Inches(0.30), size=10.5, color=SEC)
        cy += Inches(0.33)


def divider(slide, y):
    rect(slide, Inches(0.55), y, Inches(12.23), Inches(0.016), BORDER, radius=False)


def bullet_section(slide, items, x, y, w, gap=Inches(0.37)):
    cy = y
    for text, bold in items:
        if text == '':
            cy += Inches(0.10)
            continue
        prefix = '' if bold else '     ▸  '
        c = WHITE if bold else SEC
        txt(slide, prefix + text, x, cy, w, Inches(0.33),
            size=12 if bold else 11, bold=bold, color=c)
        cy += gap if bold else Inches(0.33)
    return cy


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

# Bold left stripe
rect(slide, 0, 0, Inches(0.55), H, ORANGE, radius=False)

# Top-right decorative corner block
rect(slide, Inches(10.5), 0, Inches(2.833), Inches(1.2), CARD2, radius=False)
rect(slide, Inches(10.5), 0, Inches(2.833), Inches(0.08), ORANGE, radius=False)

# Course tag
rect(slide, Inches(0.85), Inches(0.42), Inches(3.8), Inches(0.42), CARD, BLUE, 0.8)
txt(slide, 'CMPE 273  ·  Final Project  ·  Spring 2026',
    Inches(0.95), Inches(0.46), Inches(3.6), Inches(0.32),
    size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Main title
txt(slide, 'EraseGraph', Inches(0.85), Inches(1.05), Inches(9.0), Inches(2.4),
    size=82, bold=True, color=WHITE)

# Orange underline accent
rect(slide, Inches(0.85), Inches(3.12), Inches(4.8), Inches(0.07), ORANGE, radius=False)

# Tagline
txt(slide,
    'Verifiable deletion propagation across distributed systems\n'
    '— observable, traceable, and cryptographically auditable.',
    Inches(0.85), Inches(3.28), Inches(8.0), Inches(1.0),
    size=15, color=SEC)

# Team chips
team = [('Haoyuan Shan', BLUE), ('Vritika Malhotra', PURPLE),
        ('Sakshat Patil', ORANGE), ('Asim Mohammed', GREEN)]
tx, ty = Inches(0.85), Inches(4.48)
for name, color in team:
    nw = Inches(len(name) * 0.115 + 0.55)
    rect(slide, tx, ty, nw, Inches(0.48), CARD, color, 0.9)
    txt(slide, name, tx + Inches(0.14), ty + Inches(0.09), nw - Inches(0.28),
        Inches(0.32), size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx += nw + Inches(0.18)

# Right panel: architecture stack preview
rx, ry = Inches(9.5), Inches(1.1)
mini = [
    ('Frontend Dashboard',   BLUE),
    ('Backend Orchestrator', ORANGE),
    ('RabbitMQ Event Bus',   PURPLE),
    ('5 Cleanup Workers',    GREEN),
    ('Proof & Audit Chain',  AMBER),
]
for i, (label, color) in enumerate(mini):
    rect(slide, rx, ry, Inches(3.55), Inches(0.56), CARD, color, 0.9)
    txt(slide, label, rx + Inches(0.2), ry + Inches(0.13), Inches(3.15), Inches(0.32),
        size=12, bold=True, color=color)
    if i < len(mini) - 1:
        rect(slide, rx + Inches(1.6), ry + Inches(0.56), Inches(0.05),
             Inches(0.18), BORDER, radius=False)
    ry += Inches(0.74)

# Bottom tagline
txt(slide, 'GDPR / CCPA compliant  ·  Event-driven  ·  10+ microservices  ·  Kubernetes-ready',
    Inches(0.85), Inches(6.95), Inches(12.0), Inches(0.38),
    size=10, color=MUTED, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, RED)
slide_title(slide, 'The Deletion Problem')
slide_sub(slide, 'User data lives in multiple systems simultaneously — removing it from one leaves stale copies everywhere else.')

bullets = [
    ('Fragmented data stores', True),
    ('Primary DB, cache, search, analytics, and backup each hold a copy.', False),
    ('Single-system deletes are not enough', True),
    ('Removing from Postgres leaves Redis, Elasticsearch, S3, analytics untouched.', False),
    ('Compliance demands proof', True),
    ('GDPR Art.17 & CCPA require auditable evidence of propagated deletion.', False),
    ('No visibility into partial failures', True),
    ('Without tracing, a partial failure is silent — some services succeed, others fail.', False),
]
cy = Inches(1.55)
for text, bold in bullets:
    prefix = '' if bold else '     ▸  '
    c, sz = (WHITE, 12.5) if bold else (SEC, 11)
    txt(slide, prefix + text, Inches(0.55), cy, Inches(5.9), Inches(0.35),
        size=sz, bold=bold, color=c)
    cy += Inches(0.40) if bold else Inches(0.35)

image(slide, img_path('diag_problem.png'), Inches(6.6), Inches(1.35), Inches(6.45), Inches(5.85))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — SOLUTION OVERVIEW
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, BLUE)
slide_title(slide, 'Our Solution: Event-Driven Orchestration')
slide_sub(slide, 'One request triggers coordinated, parallel deletion across every data store — with a cryptographic audit trail.')

image(slide, img_path('diag_solution_flow.png'), Inches(0.35), Inches(1.5), Inches(12.63), Inches(2.55))

cards_data = [
    (ORANGE, 'Orchestrated Cleanup',
     ['One API request → 5 systems cleaned in parallel',
      'Event-driven workers — independently deployable',
      'Backend aggregates results into a final status']),
    (BLUE, 'Resilient by Default',
     ['Retry queues + dead-letter exchange (DLQ)',
      'Circuit breakers prevent cascade failures',
      'Idempotency guards prevent double-processing']),
    (GREEN, 'Cryptographic Audit Trail',
     ['Every event SHA-256 hash-chained',
      'Tamper-evident — modify any record, chain breaks',
      'Exportable JSON proof + one-click verification']),
]
cx = Inches(0.35)
for color, title, blist in cards_data:
    card(slide, cx, Inches(4.25), Inches(4.17), Inches(2.95), color, title, blist)
    cx += Inches(4.32)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — SYSTEM ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, PURPLE)
slide_title(slide, 'System Architecture')
slide_sub(slide, 'Frontend · Backend · RabbitMQ · 5 Cleanup Workers · Proof Service · PostgreSQL · Redis · S3')

image(slide, img_path('diag_architecture.png'), Inches(0.25), Inches(1.38), Inches(12.83), Inches(5.85))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — SEQUENCE DIAGRAM  (NEW)
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, TEAL)
slide_title(slide, 'Request Sequence — Step by Step', color=WHITE)
slide_sub(slide, 'How a single deletion request flows through every service layer — from REST call to SSE completion event.')

image(slide, img_path('diag_sequence.png'), Inches(0.25), Inches(1.38), Inches(12.83), Inches(5.85))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — WORKFLOW / STATE MACHINE
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, AMBER)
slide_title(slide, 'Request Lifecycle & State Machine')
slide_sub(slide, 'Every deletion moves through a well-defined state machine with step-level visibility per service.')

image(slide, img_path('diag_workflow.png'), Inches(0.35), Inches(1.38), Inches(12.63), Inches(4.4))

rect(slide, Inches(0.35), Inches(5.95), Inches(6.1), Inches(1.3), CARD, GREEN, 0.9)
txt(slide, 'Step states per service', Inches(0.55), Inches(6.08), Inches(5.7), Inches(0.32),
    size=11.5, bold=True, color=GREEN)
txt(slide, 'Each worker independently reports PENDING → RUNNING → SUCCEEDED / FAILED. '
    'The orchestrator aggregates into a final request-level status.',
    Inches(0.55), Inches(6.42), Inches(5.7), Inches(0.65), size=10.5, color=SEC)

rect(slide, Inches(6.65), Inches(5.95), Inches(6.32), Inches(1.3), CARD, BLUE, 0.9)
txt(slide, 'PARTIAL_COMPLETED is intentional', Inches(6.85), Inches(6.08), Inches(5.9), Inches(0.32),
    size=11.5, bold=True, color=BLUE)
txt(slide, 'Non-critical services (analytics, backup) failing does not block primary deletion. '
    'Auditors see exactly which systems were cleaned — more informative than binary pass/fail.',
    Inches(6.85), Inches(6.42), Inches(5.9), Inches(0.65), size=10.5, color=SEC)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — RELIABILITY
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, ORANGE)
slide_title(slide, 'Reliability Design')
slide_sub(slide, 'Three layered mechanisms ensure every deletion eventually completes — or is explicitly accounted for.')

image(slide, img_path('diag_reliability.png'), Inches(0.25), Inches(1.38), Inches(12.83), Inches(5.85))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — PROOF & AUDITABILITY
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, GREEN)
slide_title(slide, 'Proof & Auditability')
slide_sub(slide, 'Every deletion event is hash-chained. Modify any record and verification fails — cryptographically guaranteed.')

proof_bullets = [
    ('Hash-chained proof events', True),
    ('Each event stores SHA-256(prev_hash + type + payload + ts).', False),
    ('Any tampering breaks the chain from that point onwards.', False),
    ('', False),
    ('Advisory lock during writes', True),
    ('pg_advisory_xact_lock ensures sequential writes', False),
    ('under concurrent load — no race conditions.', False),
    ('', False),
    ('Export & one-click verification', True),
    ('Full chain exported as JSON. Verify reports', False),
    ('VERIFIED or TAMPERED per-record.', False),
    ('', False),
    ('Compliance-ready', True),
    ('Satisfies GDPR Art.17 & CCPA deletion audit requirements.', False),
]
cy = Inches(1.55)
for text, bold in proof_bullets:
    if text == '':
        cy += Inches(0.08)
        continue
    prefix = '' if bold else '     ▸  '
    c, sz = (WHITE, 12.5) if bold else (SEC, 10.5)
    txt(slide, prefix + text, Inches(0.55), cy, Inches(5.6), Inches(0.32),
        size=sz, bold=bold, color=c)
    cy += Inches(0.36) if bold else Inches(0.32)

image(slide, img_path('diag_proof_chain.png'), Inches(6.4), Inches(1.38), Inches(6.65), Inches(5.85))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — CONSISTENCY & TRADEOFFS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, AMBER)
slide_title(slide, 'Consistency & Tradeoffs')
slide_sub(slide, 'Deliberate tradeoffs per data store. Not all systems require the same consistency guarantee.')

image(slide, img_path('diag_consistency.png'), Inches(0.35), Inches(1.38), Inches(12.63), Inches(4.2))

rect(slide, Inches(0.35), Inches(5.75), Inches(6.1), Inches(1.5), CARD, ORANGE, 0.9)
txt(slide, 'Why PARTIAL_COMPLETED is good design', Inches(0.55), Inches(5.88), Inches(5.7), Inches(0.32),
    size=11.5, bold=True, color=ORANGE)
txt(slide, 'Rather than failing entirely when analytics is slow, we record partial success. '
    'Auditors see exactly which systems were cleaned — more informative than binary pass/fail.',
    Inches(0.55), Inches(6.24), Inches(5.7), Inches(0.80), size=10.5, color=SEC)

rect(slide, Inches(6.65), Inches(5.75), Inches(6.32), Inches(1.5), CARD, BLUE, 0.9)
txt(slide, 'CAP Theorem position', Inches(6.85), Inches(5.88), Inches(5.9), Inches(0.32),
    size=11.5, bold=True, color=BLUE)
txt(slide, 'We favour AP (Availability + Partition Tolerance) for the event-driven path. '
    'The cryptographic proof chain provides Consistency through eventual convergence and tamper detection.',
    Inches(6.85), Inches(6.24), Inches(5.9), Inches(0.80), size=10.5, color=SEC)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — OBSERVABILITY
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, TEAL)
slide_title(slide, 'Observability & Operations')
slide_sub(slide, 'Every request, event, and failure is traced, measured, and surfaced in real time.')

obs_items = [
    ('OpenTelemetry + Jaeger', True),
    ('Distributed traces span all service boundaries.', False),
    ('See exactly where latency and failures occur.', False),
    ('', False),
    ('Prometheus + Grafana', True),
    ('Throughput, step duration, retry counts, CB state', False),
    ('exported and visualised on live dashboards.', False),
    ('', False),
    ('Admin Dashboard', True),
    ('Real-time health, circuit breaker state per service,', False),
    ('SLA violations with timestamps, DLQ message count.', False),
    ('', False),
    ('Real-time SSE Streaming', True),
    ('Frontend receives live status updates — no polling.', False),
]
cy = Inches(1.55)
for text, bold in obs_items:
    if text == '':
        cy += Inches(0.08)
        continue
    prefix = '' if bold else '     ▸  '
    c, sz = (WHITE, 12.5) if bold else (SEC, 10.5)
    txt(slide, prefix + text, Inches(0.55), cy, Inches(6.1), Inches(0.32),
        size=sz, bold=bold, color=c)
    cy += Inches(0.36) if bold else Inches(0.32)

rx, ry = Inches(7.1), Inches(1.55)
rect(slide, rx, ry, Inches(5.95), Inches(1.7), CARD, TEAL, 1.0)
txt(slide, 'Observability Stack', rx + Inches(0.22), ry + Inches(0.17), Inches(5.5), Inches(0.32),
    size=12, bold=True, color=TEAL)
tools = [('OpenTelemetry', TEAL), ('Jaeger', BLUE), ('Prometheus', ORANGE),
         ('Grafana', GREEN), ('SSE Streaming', PURPLE)]
px2, py2 = rx + Inches(0.22), ry + Inches(0.60)
for label, color in tools:
    pw = pill(slide, label, px2, py2, color)
    px2 += pw + Inches(0.14)
    if px2 > rx + Inches(5.6):
        px2 = rx + Inches(0.22)
        py2 += Inches(0.42)

rect(slide, rx, ry + Inches(1.90), Inches(5.95), Inches(2.8), CARD, GREEN, 1.0)
txt(slide, 'Admin Panel surfaces', rx + Inches(0.22), ry + Inches(2.05), Inches(5.5), Inches(0.32),
    size=12, bold=True, color=GREEN)
admin_items = [
    'Per-service health (healthy / degraded / down)',
    'Circuit breaker state — CLOSED / OPEN / HALF-OPEN',
    'SLA violations with timestamps and service name',
    'DLQ message count + one-click replay trigger',
    'Cryptographic attestation per deletion request',
]
ay = ry + Inches(2.46)
for item in admin_items:
    txt(slide, '  ▸  ' + item, rx + Inches(0.22), ay, Inches(5.5), Inches(0.30),
        size=10.5, color=SEC)
    ay += Inches(0.33)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — DEPLOYMENT
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, BLUE)
slide_title(slide, 'Kubernetes-Ready, Cloud-Deployed')
slide_sub(slide, 'Every service containerised — local via Docker Compose, production-grade on Kubernetes / GKE.')

dep_items = [
    ('Docker Compose (local dev)', True),
    ('Single command spins up all 12+ services,', False),
    ('Postgres, Redis, RabbitMQ, Jaeger, Prometheus.', False),
    ('', False),
    ('Kubernetes (production)', True),
    ('Deployment manifests, Services, ConfigMaps for every', False),
    ('component. Workers scale horizontally.', False),
    ('', False),
    ('Cloud Deployment (GKE)', True),
    ('Live on cloud cluster. All services verified running', False),
    ('end-to-end in a production-like environment.', False),
    ('', False),
    ('Environment isolation', True),
    ('Secrets via .env.example — zero credentials committed.', False),
]
cy = Inches(1.55)
for text, bold in dep_items:
    if text == '':
        cy += Inches(0.08)
        continue
    prefix = '' if bold else '     ▸  '
    c, sz = (WHITE, 12.5) if bold else (SEC, 10.5)
    txt(slide, prefix + text, Inches(0.55), cy, Inches(6.3), Inches(0.32),
        size=sz, bold=bold, color=c)
    cy += Inches(0.36) if bold else Inches(0.32)

rx = Inches(7.3)
rect(slide, rx, Inches(1.55), Inches(5.7), Inches(2.9), CARD, BLUE, 1.0)
txt(slide, 'Services deployed', rx + Inches(0.22), Inches(1.70), Inches(5.26), Inches(0.32),
    size=12, bold=True, color=BLUE)
services = ['Backend API', 'Frontend', 'Cache Cleanup', 'Search Cleanup',
            'Analytics Cleanup', 'Backup Service', 'Proof Service',
            'RabbitMQ', 'PostgreSQL', 'Redis', 'Jaeger', 'Prometheus']
sx2, sy2 = rx + Inches(0.22), Inches(2.12)
for svc in services:
    pw = pill(slide, svc, sx2, sy2, MUTED, size=9)
    sx2 += pw + Inches(0.12)
    if sx2 > rx + Inches(5.4):
        sx2 = rx + Inches(0.22)
        sy2 += Inches(0.40)

metrics = [('65', 'Unit Tests', BLUE), ('15', 'Test Suites', PURPLE),
           ('5', 'Data Stores', ORANGE), ('100%', 'Passing', GREEN)]
mx = rx
for val, label, color in metrics:
    rect(slide, mx, Inches(4.65), Inches(1.3), Inches(1.4), CARD, color, 0.9)
    txt(slide, val, mx, Inches(4.78), Inches(1.3), Inches(0.72),
        size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label, mx, Inches(5.4), Inches(1.3), Inches(0.30),
        size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    mx += Inches(1.45)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — KEY RESULTS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, GREEN)
slide_title(slide, 'Key Results')
slide_sub(slide, 'What we built and demonstrated across the full semester.')

big_metrics = [
    ('65',  'Unit Tests Passing', BLUE),
    ('15',  'Test Suites',        PURPLE),
    ('5',   'Data Stores',        ORANGE),
    ('3',   'Reliability Layers', GREEN),
    ('10+', 'Microservices',      AMBER),
]
mx = Inches(0.35)
mw = Inches(2.5)
for val, label, color in big_metrics:
    rect(slide, mx, Inches(1.58), mw, Inches(1.85), CARD, color, 1.0)
    txt(slide, val, mx, Inches(1.74), mw, Inches(1.05),
        size=52, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label, mx + Inches(0.1), Inches(2.78), mw - Inches(0.2), Inches(0.42),
        size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    mx += mw + Inches(0.10)

result_cards = [
    (ORANGE, 'Complete Distributed System',
     ['10+ microservices coordinating deletion across 5 data store types.',
      'Event-driven, independently deployable workers.',
      'Full lifecycle: request intake → audit export.']),
    (BLUE, 'Verifiable Proof Chain',
     ['Cryptographic hash chain for every deletion request.',
      'Tamper detection works at the single-record level.',
      'One-click verify: VERIFIED / TAMPERED result.']),
    (GREEN, 'Failure-Safe by Design',
     ['Retry + DLQ + idempotency + circuit breaker.',
      'No deletion lost silently, even under partitions.',
      'DLQ replay ensures ops can recover any failed step.']),
]
cx = Inches(0.35)
for color, title, blist in result_cards:
    card(slide, cx, Inches(3.68), Inches(4.17), Inches(3.4), color, title, blist)
    cx += Inches(4.32)

txt(slide,
    'EraseGraph demonstrates that compliance-grade deletion propagation is achievable with event-driven '
    'orchestration, reliability engineering, and cryptographic auditability.',
    Inches(0.55), Inches(7.10), Inches(12.23), Inches(0.34),
    size=10.5, color=MUTED)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — CHALLENGES & LESSONS LEARNED  (NEW)
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, RED)
slide_title(slide, 'Challenges & Lessons Learned')
slide_sub(slide, 'Building a reliable distributed system is harder than it looks — here is what we learned.')

challenges = [
    ('Distributed idempotency', ORANGE,
     ['Ensuring exactly-once processing across async workers',
      'Race conditions on duplicate events required DB-level unique constraints',
      'Clock skew makes timestamp-only deduplication unreliable']),
    ('Hash chain concurrency', GREEN,
     ['Concurrent inserts can corrupt prev_hash references',
      'Solved with pg_advisory_xact_lock for sequential writes',
      'Verified correct under load with integration tests']),
    ('Circuit breaker persistence', PURPLE,
     ['State must survive worker restarts — stored in Redis',
      'Half-open probe logic required careful timeout tuning',
      'Unexpected: CB state affects DLQ routing decisions']),
    ('Trace context across async', TEAL,
     ['OpenTelemetry span context lost across RabbitMQ boundaries',
      'Required manual W3C trace-context header propagation',
      'Jaeger now shows full end-to-end trace per request']),
]

cx = Inches(0.35)
for title, color, bullets in challenges:
    cw = Inches(3.12)
    ch = Inches(5.1)
    rect(slide, cx, Inches(1.55), cw, ch, CARD, color, 1.0)
    txt(slide, title, cx + Inches(0.18), Inches(1.70), cw - Inches(0.36),
        Inches(0.34), size=11.5, bold=True, color=color)
    by = Inches(2.18)
    for b in bullets:
        txt(slide, '  ▸  ' + b, cx + Inches(0.18), by, cw - Inches(0.36),
            Inches(0.55), size=10, color=SEC, wrap=True)
        by += Inches(0.62)
    cx += cw + Inches(0.12)

# Key lesson banner at bottom
rect(slide, Inches(0.35), Inches(6.80), Inches(12.63), Inches(0.50), CARD, AMBER, 1.0)
txt(slide,
    '  Key lesson: Event-driven fanout beats synchronous RPC for deletion — lower coupling, '
    'better scalability, explicit failure visibility.',
    Inches(0.52), Inches(6.88), Inches(12.3), Inches(0.36),
    size=11, bold=False, color=AMBER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — FUTURE WORK  (NEW)
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, PURPLE)
slide_title(slide, 'Future Work')
slide_sub(slide, 'Three directions to take EraseGraph from a proof-of-concept to a production-grade platform.')

future = [
    (BLUE,   'Scalability',
     'Replace RabbitMQ with Apache Kafka',
     ['Persistent, replayable event log for replay-on-demand',
      'Consumer group auto-scaling with KEDA on Kubernetes',
      'Partition deletion requests by user_id for parallelism',
      'Exactly-once delivery with Kafka transactions']),
    (GREEN,  'Compliance Automation',
     'GDPR / CCPA deletion certificate generation',
     ['Automated PDF deletion certificate with proof hash',
      'SLA violation reports emailed to data officers',
      'Multi-region data residency tagging per request',
      'Webhook notifications to data subjects on completion']),
    (PURPLE, 'Intelligence & Resilience',
     'Observability-driven optimisation',
     ['ML anomaly detection on deletion latency patterns',
      'GraphQL audit query API for flexible reporting',
      'Chaos engineering suite (fault injection, partition tests)',
      'Auto-remediation: trigger replay on SLA breach detection']),
]

cx = Inches(0.35)
for color, title, subtitle, bullets in future:
    cw, ch = Inches(4.12), Inches(5.3)
    rect(slide, cx, Inches(1.55), cw, ch, CARD, color, 1.0)

    # Colour top bar inside card
    rect(slide, cx, Inches(1.55), cw, Inches(0.06), color, radius=False)

    txt(slide, title, cx + Inches(0.2), Inches(1.70), cw - Inches(0.4),
        Inches(0.36), size=13, bold=True, color=color)
    txt(slide, subtitle, cx + Inches(0.2), Inches(2.10), cw - Inches(0.4),
        Inches(0.30), size=10.5, color=MUTED, italic=True)

    by = Inches(2.52)
    for b in bullets:
        txt(slide, '  ▸  ' + b, cx + Inches(0.2), by, cw - Inches(0.4),
            Inches(0.52), size=10.5, color=SEC, wrap=True)
        by += Inches(0.60)
    cx += cw + Inches(0.18)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — TEAM CONTRIBUTIONS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
accent_bar(slide, PURPLE)
slide_title(slide, 'Team Contributions')
slide_sub(slide, 'Responsibility breakdown — who owned what and the impact delivered.')

contrib = [
    ('Haoyuan Shan',    BLUE,
     'GitHub consolidation · K8s / cloud deployment · backend integration',
     'K8s manifests, cloud deployment, proof-chain optimisation, demo data',
     'Turned the project into a cohesive, demo-ready deliverable'),
    ('Vritika Malhotra', PURPLE,
     'Retry · DLQ · idempotency · circuit breaker · reliability docs',
     'Circuit breaker service, DLQ replay, idempotency guard, unit tests',
     'Improved failure recovery & robustness of distributed deletion'),
    ('Sakshat Patil',   ORANGE,
     'Frontend dashboard · workflow UX · real-time status experience',
     'React pages, SSE streaming UI, bulk upload, admin front-end',
     'Improved visualisation quality and demo clarity'),
    ('Asim Mohammed',   GREEN,
     'Consistency semantics · backend / admin workflow · state model',
     'SLA monitor, SSE backend, deletion-request service, state design',
     'Improved state-model clarity and operational reliability'),
]

col_x = [Inches(0.35), Inches(2.2), Inches(5.55), Inches(9.0)]
col_w = [Inches(1.75), Inches(3.25), Inches(3.35), Inches(4.0)]
headers = ['Member', 'Owned Areas', 'Key Deliverables', 'Impact']

rect(slide, Inches(0.35), Inches(1.56), Inches(12.63), Inches(0.40), CARD, BORDER, 0.6)
for i, h in enumerate(headers):
    txt(slide, h, col_x[i] + Inches(0.1), Inches(1.62), col_w[i], Inches(0.28),
        size=9.5, bold=True, color=MUTED)

ry = Inches(1.96)
rh = Inches(1.22)
for name, color, areas, deliverables, impact in contrib:
    rect(slide, Inches(0.35), ry, Inches(12.63), rh - Inches(0.05), CARD, color, 0.5)
    for i, val in enumerate([name, areas, deliverables, impact]):
        bold = (i == 0)
        c = color if bold else SEC
        txt(slide, val, col_x[i] + Inches(0.1), ry + Inches(0.12),
            col_w[i] - Inches(0.12), rh - Inches(0.20),
            size=11 if bold else 9.5, bold=bold, color=c)
    ry += rh

txt(slide, 'R = Responsible   A = Accountable   C = Consulted   I = Informed',
    Inches(0.55), Inches(7.02), Inches(12.0), Inches(0.28),
    size=8.5, color=MUTED, align=PP_ALIGN.RIGHT, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Q&A
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
rect(slide, 0, 0, Inches(0.55), H, ORANGE, radius=False)
rect(slide, 0, 0, W, Inches(0.07), ORANGE, radius=False)

txt(slide, 'Thank you.', Inches(0.9), Inches(1.2), Inches(11.5), Inches(2.2),
    size=88, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

rect(slide, Inches(0.9), Inches(3.25), Inches(5.0), Inches(0.07), ORANGE, radius=False)

txt(slide, 'Questions?', Inches(0.9), Inches(3.4), Inches(11.5), Inches(1.5),
    size=72, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

txt(slide, 'EraseGraph  —  Verifiable deletion propagation  ·  CMPE 273 Spring 2026',
    Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.42),
    size=13, color=SEC, align=PP_ALIGN.LEFT)

tags = [('Event-Driven', BLUE), ('Cryptographic Proof', GREEN),
        ('Reliability First', ORANGE), ('GDPR / CCPA Ready', PURPLE),
        ('10+ Microservices', AMBER), ('Kubernetes', TEAL)]
px2, py2 = Inches(0.9), Inches(5.55)
for label, color in tags:
    pw = pill(slide, label, px2, py2, color, size=10.5)
    px2 += pw + Inches(0.22)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out_file = os.path.join(OUT, 'erasegraph-slides.pptx')
prs.save(out_file)
n = len(prs.slides)
size = os.path.getsize(out_file) // 1024
print(f'Saved: {out_file}')
print(f'{n} slides  ·  {size} KB')
