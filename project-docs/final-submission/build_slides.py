"""
EraseGraph — CMPE 273 Final Presentation
Generates erasegraph-slides.pptx matching the frontend dark navy + orange/blue theme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Canvas: widescreen 16:9 ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Colour palette (matches index.css) ───────────────────────────────────────
BG_DARK    = RGBColor(0x08, 0x11, 0x1f)   # --bg-top
SURFACE    = RGBColor(0x0b, 0x12, 0x21)   # --surface-card
SURFACE_EL = RGBColor(0x0f, 0x18, 0x2a)   # --surface-elevated
BORDER     = RGBColor(0x95, 0xb1, 0xe5)   # --border-soft (at 14% alpha — used solid at low alpha)
TEXT_PRI   = RGBColor(0xee, 0xf4, 0xff)   # --text-primary
TEXT_SEC   = RGBColor(0x9e, 0xaf, 0xcf)   # --text-secondary
TEXT_MUT   = RGBColor(0x71, 0x80, 0xa0)   # --text-muted
ORANGE     = RGBColor(0xff, 0x7a, 0x45)   # --accent-strong
BLUE       = RGBColor(0x4e, 0x8d, 0xff)   # --accent-secondary
GREEN      = RGBColor(0x7c, 0xf0, 0xbc)   # --accent-green
PURPLE     = RGBColor(0xc0, 0x84, 0xfc)   # --accent-purple
AMBER      = RGBColor(0xfc, 0xd3, 0x4d)   # --amber

# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    """Add a slide with a blank layout and dark background."""
    blank_layout = prs.slide_layouts[6]   # truly blank
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide


def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width_pt=0.75):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width_pt)
    else:
        shape.line.fill.background()
    # Round corners via XML
    sp = shape.element
    prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        else:
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 20000')   # ~corner radius
    return shape


def add_text_box(slide, text, x, y, w, h, font_size=18, bold=False,
                 color=TEXT_PRI, align=PP_ALIGN.LEFT, font_name="Inter", italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def heading(slide, text, x=Inches(0.55), y=Inches(0.28), w=Inches(12.2), font_size=36):
    return add_text_box(slide, text, x, y, w, Inches(1), font_size=font_size,
                        bold=True, color=TEXT_PRI, font_name="Calibri")


def subheading(slide, text, x=Inches(0.55), y=Inches(1.12), w=Inches(12.2)):
    return add_text_box(slide, text, x, y, w, Inches(0.5), font_size=14,
                        color=TEXT_SEC, font_name="Calibri")


def label_pill(slide, text, x, y, w=Inches(2), color=BLUE):
    box = add_rect(slide, x, y, w, Inches(0.28), SURFACE_EL, color, 0.5)
    add_text_box(slide, text, x, y + Inches(0.04), w, Inches(0.24),
                 font_size=9, bold=True, color=color,
                 align=PP_ALIGN.CENTER, font_name="Calibri")
    return box


def card(slide, x, y, w, h, accent_color=None):
    c = SURFACE_EL if accent_color is None else SURFACE_EL
    bdr = BORDER if accent_color is None else accent_color
    return add_rect(slide, x, y, w, h, c, bdr, 0.6)


def card_with_title(slide, title, body_lines, x, y, w, h,
                    icon="", accent_color=None, title_color=TEXT_PRI):
    card(slide, x, y, w, h, accent_color)
    ty = y + Inches(0.15)
    if icon:
        add_text_box(slide, icon, x + Inches(0.18), ty, Inches(0.4), Inches(0.35),
                     font_size=16, align=PP_ALIGN.LEFT, font_name="Segoe UI Emoji")
        tx = x + Inches(0.6)
    else:
        tx = x + Inches(0.18)
    add_text_box(slide, title, tx, ty, w - Inches(0.8), Inches(0.3),
                 font_size=12, bold=True, color=title_color or (accent_color if accent_color else TEXT_PRI), font_name="Calibri")
    body_y = ty + Inches(0.3)
    for line in body_lines:
        add_text_box(slide, line, x + Inches(0.18), body_y, w - Inches(0.36), Inches(0.28),
                     font_size=10, color=TEXT_SEC, font_name="Calibri")
        body_y += Inches(0.27)


def bullet_card(slide, title, body, x, y, w, icon="", accent=None):
    h = Inches(0.85)
    card(slide, x, y, w, h, accent)
    if icon:
        add_text_box(slide, icon, x + Inches(0.14), y + Inches(0.12), Inches(0.35), Inches(0.4),
                     font_size=18, font_name="Segoe UI Emoji")
        tx = x + Inches(0.55)
        tw = w - Inches(0.7)
    else:
        tx = x + Inches(0.18)
        tw = w - Inches(0.36)
    add_text_box(slide, title, tx, y + Inches(0.10), tw, Inches(0.28),
                 font_size=11, bold=True, color=TEXT_PRI, font_name="Calibri")
    add_text_box(slide, body, tx, y + Inches(0.38), tw, Inches(0.40),
                 font_size=9.5, color=TEXT_SEC, font_name="Calibri")


def section_label(slide, text, x=Inches(0.55), y=Inches(0.12)):
    label_pill(slide, text, x, y, Inches(2.2), BLUE)


def arrow_h(slide, x1, y, x2, color=BORDER):
    """Horizontal arrow from (x1,y) to (x2,y)."""
    from pptx.util import Inches
    connector = slide.shapes.add_connector(1, x1, y, x2, y)
    connector.line.color.rgb = color
    connector.line.width = Pt(1.2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 · TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)

# Orange accent strip (top)
add_rect(slide, 0, 0, W, Inches(0.06), ORANGE)

# Course badge
label_pill(slide, "CMPE 273  ·  Final Project  ·  Spring 2026",
           Inches(0.55), Inches(0.22), Inches(3.4), BLUE)

# Title
add_text_box(slide, "EraseGraph", Inches(0.55), Inches(0.72), Inches(7.5), Inches(1.9),
             font_size=72, bold=True, color=TEXT_PRI, font_name="Calibri")

# Gradient effect on "Graph" — second text box overlaid in orange
add_text_box(slide, "           Graph", Inches(0.55), Inches(0.72), Inches(7.5), Inches(1.9),
             font_size=72, bold=True, color=ORANGE, font_name="Calibri")

# Tagline
add_text_box(slide,
             "Verifiable deletion propagation in distributed systems —\nobservable, traceable, and auditable end-to-end.",
             Inches(0.55), Inches(2.45), Inches(7.5), Inches(0.9),
             font_size=15, color=TEXT_SEC, font_name="Calibri")

# Team members
team = ["Haoyuan Shan", "Vritika Malhotra", "Sakshat Patil", "Asim Mohammed"]
for i, name in enumerate(team):
    bx = Inches(0.55 + i * 2.05)
    add_rect(slide, bx, Inches(3.6), Inches(1.95), Inches(0.42), SURFACE_EL, BORDER, 0.5)
    add_text_box(slide, name, bx + Inches(0.1), Inches(3.63), Inches(1.75), Inches(0.35),
                 font_size=11, bold=True, color=TEXT_PRI, font_name="Calibri", align=PP_ALIGN.CENTER)

# Right side mini-arch
mini_items = [
    ("Frontend Dashboard", BLUE),
    ("Backend Orchestrator", ORANGE),
    ("RabbitMQ Event Bus", PURPLE),
    ("5 Cleanup Workers", GREEN),
    ("Proof & Audit Chain", AMBER),
]
rx = Inches(9.1)
ry = Inches(0.9)
for label, color in mini_items:
    add_rect(slide, rx, ry, Inches(3.9), Inches(0.46), SURFACE_EL, color, 0.7)
    add_text_box(slide, label, rx + Inches(0.14), ry + Inches(0.09), Inches(3.6), Inches(0.3),
                 font_size=11, bold=True, color=color, font_name="Calibri")
    ry += Inches(0.58)
    if ry < Inches(0.9) + Inches(0.58) * 4.5:
        # connector dot
        add_rect(slide, rx + Inches(1.8), ry - Inches(0.1), Inches(0.04), Inches(0.2),
                 BORDER, None)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 · PROBLEM & MOTIVATION
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), ORANGE)
section_label(slide, "Problem & Motivation")
heading(slide, "Deletion is harder than it looks", y=Inches(0.48))
subheading(slide, "User data lives in many places. A single-service delete leaves stale copies silently behind.")

bullets = [
    ("🗄️", "Fragmented data stores",
     "Primary DB, cache, search index, analytics warehouse, and backup — each holds a copy."),
    ("🚫", "Single-system deletes aren't enough",
     "Removing from Postgres leaves it live in Redis, Elasticsearch, S3, and analytics."),
    ("📋", "Compliance demands proof",
     "GDPR / CCPA require auditable evidence that deletion propagated — not just a log entry."),
    ("🔍", "No visibility into failures",
     "Without tracing, a partial failure is invisible — some services succeed, others silently fail."),
]
by = Inches(1.7)
for icon, title, body in bullets:
    bullet_card(slide, title, body, Inches(0.55), by, Inches(6.1), icon=icon, accent=ORANGE)
    by += Inches(0.97)

# Data scatter diagram (right side)
stores = [
    (Inches(8.1),  Inches(1.6),  "🗄️  PostgreSQL",  "Primary DB",       ORANGE),
    (Inches(10.3), Inches(1.6),  "⚡  Redis",        "Cache Layer",       BLUE),
    (Inches(8.1),  Inches(2.75), "🔍  Elasticsearch","Search Index",      PURPLE),
    (Inches(10.3), Inches(2.75), "📊  Analytics",    "Data Warehouse",    GREEN),
    (Inches(9.2),  Inches(3.9),  "💾  Backup",       "S3 / Blob Store",   AMBER),
]
# User node
add_rect(slide, Inches(9.45), Inches(0.5), Inches(1.7), Inches(0.7), SURFACE_EL, BLUE, 0.8)
add_text_box(slide, "👤  User", Inches(9.45), Inches(0.62), Inches(1.7), Inches(0.35),
             font_size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Segoe UI Emoji")

for sx, sy, label, sub, color in stores:
    add_rect(slide, sx, sy, Inches(1.95), Inches(0.68), SURFACE_EL, color, 0.6)
    add_text_box(slide, label, sx + Inches(0.1), sy + Inches(0.06), Inches(1.75), Inches(0.3),
                 font_size=9.5, bold=True, color=color, font_name="Calibri")
    add_text_box(slide, sub, sx + Inches(0.1), sy + Inches(0.36), Inches(1.75), Inches(0.25),
                 font_size=8.5, color=TEXT_MUT, font_name="Calibri")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 · SOLUTION OVERVIEW
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), BLUE)
section_label(slide, "Solution Overview")
heading(slide, "Event-driven orchestration with proof", y=Inches(0.48))
subheading(slide, "EraseGraph coordinates deletion across every data store and produces an immutable audit trail — automatically.")

# Flow nodes
flow_nodes = [
    ("Deletion\nRequest",    BLUE,   Inches(0.35)),
    ("Backend\nOrchestrator",ORANGE, Inches(2.1)),
    ("RabbitMQ\nEvent Bus",  PURPLE, Inches(3.85)),
    ("Cleanup\nWorkers ×5",  GREEN,  Inches(5.6)),
    ("Proof\nAudit Chain",   AMBER,  Inches(7.35)),
    ("Admin\nDashboard",     BLUE,   Inches(9.1)),
]
nw, nh = Inches(1.55), Inches(0.72)
ny = Inches(1.75)
for label, color, nx in flow_nodes:
    add_rect(slide, nx, ny, nw, nh, SURFACE_EL, color, 0.8)
    add_text_box(slide, label, nx + Inches(0.08), ny + Inches(0.1), nw - Inches(0.16), nh - Inches(0.1),
                 font_size=10, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Calibri")

# Arrows between nodes
arrow_labels = ["REST", "AMQP", "Fan-out", "Events", "View"]
for i in range(len(flow_nodes) - 1):
    ax = flow_nodes[i][2] + nw + Inches(0.02)
    ay = ny + nh / 2
    ax2 = flow_nodes[i+1][2] - Inches(0.02)
    arrow_h(slide, ax, ay, ax2)
    mid = ax + (ax2 - ax) / 2
    add_text_box(slide, arrow_labels[i], mid - Inches(0.25), ny - Inches(0.28), Inches(0.5), Inches(0.22),
                 font_size=7.5, color=TEXT_MUT, align=PP_ALIGN.CENTER, font_name="Calibri")

# Three cards
cards3 = [
    ("🎯", "Orchestrated Multi-Service Cleanup",
     ["A single request triggers coordinated deletion", "across DB, cache, search, analytics, and backup."],
     ORANGE, Inches(0.35)),
    ("🔁", "Resilient by Default",
     ["Retry queues, DLQ, circuit breakers, and", "idempotency guards make every step failure-safe."],
     BLUE, Inches(4.6)),
    ("🔐", "Cryptographic Proof Chain",
     ["Every deletion event is hash-chained and signed —", "tamper-evident audit trail, always verifiable."],
     GREEN, Inches(8.85)),
]
for icon, title, body, color, cx in cards3:
    card_with_title(slide, title, body, cx, Inches(2.85), Inches(4.15), Inches(1.8),
                    icon=icon, accent_color=color, title_color=color)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 · SYSTEM ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), PURPLE)
section_label(slide, "System Architecture")
heading(slide, "Full-stack distributed system", y=Inches(0.48))

# Layer 0: Frontend
add_rect(slide, Inches(3.9), Inches(1.35), Inches(5.5), Inches(0.62), SURFACE_EL, BLUE, 0.8)
add_text_box(slide, "Frontend Dashboard (React)",
             Inches(3.9), Inches(1.44), Inches(5.5), Inches(0.28),
             font_size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text_box(slide, "Status · Proof · Bulk Upload · Admin",
             Inches(3.9), Inches(1.72), Inches(5.5), Inches(0.22),
             font_size=9, color=TEXT_MUT, align=PP_ALIGN.CENTER, font_name="Calibri")

# Arrow: Frontend → Backend
add_rect(slide, Inches(6.45), Inches(1.97), Inches(0.42), Inches(0.5), SURFACE, BORDER, 0.3)
add_text_box(slide, "REST / SSE", Inches(6.55), Inches(2.0), Inches(0.8), Inches(0.2),
             font_size=7, color=TEXT_MUT, font_name="Calibri")

# Layer 1: Backend
add_rect(slide, Inches(2.5), Inches(2.47), Inches(8.3), Inches(0.68), SURFACE_EL, ORANGE, 0.8)
add_text_box(slide, "Backend Orchestrator (NestJS)",
             Inches(2.5), Inches(2.56), Inches(8.3), Inches(0.28),
             font_size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text_box(slide, "Workflow Coordination · Auth · Retry Logic · SLA Monitor · Proof Chain",
             Inches(2.5), Inches(2.82), Inches(8.3), Inches(0.22),
             font_size=9, color=TEXT_MUT, align=PP_ALIGN.CENTER, font_name="Calibri")

# Layer 2: RabbitMQ + Observability
add_rect(slide, Inches(0.35), Inches(3.42), Inches(3.2), Inches(0.62), SURFACE_EL, GREEN, 0.7)
add_text_box(slide, "Observability Stack",
             Inches(0.35), Inches(3.5), Inches(3.2), Inches(0.26),
             font_size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text_box(slide, "OpenTelemetry · Jaeger · Prometheus · Grafana",
             Inches(0.35), Inches(3.74), Inches(3.2), Inches(0.22),
             font_size=8, color=TEXT_MUT, align=PP_ALIGN.CENTER, font_name="Calibri")

add_rect(slide, Inches(6.8), Inches(3.42), Inches(3.5), Inches(0.62), SURFACE_EL, PURPLE, 0.8)
add_text_box(slide, "RabbitMQ",
             Inches(6.8), Inches(3.5), Inches(3.5), Inches(0.26),
             font_size=10, bold=True, color=PURPLE, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text_box(slide, "Topic Exchange · DLQ · Retry Queue",
             Inches(6.8), Inches(3.74), Inches(3.5), Inches(0.22),
             font_size=8, color=TEXT_MUT, align=PP_ALIGN.CENTER, font_name="Calibri")

# Layer 3: Workers
workers = [
    ("Cache\nCleanup",     ORANGE, Inches(0.2)),
    ("Search\nCleanup",    BLUE,   Inches(2.65)),
    ("Analytics\nCleanup", GREEN,  Inches(5.1)),
    ("Backup\nService",    AMBER,  Inches(7.55)),
    ("Proof\nService",     PURPLE, Inches(10.0)),
]
for label, color, wx in workers:
    add_rect(slide, wx, Inches(4.32), Inches(2.2), Inches(0.7), SURFACE_EL, color, 0.7)
    add_text_box(slide, label, wx + Inches(0.1), Inches(4.42), Inches(2.0), Inches(0.52),
                 font_size=10, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Calibri")

# Layer 4: Storage
add_rect(slide, Inches(0.2), Inches(5.3), Inches(4.1), Inches(0.55), SURFACE_EL, ORANGE, 0.6)
add_text_box(slide, "🗄️  PostgreSQL   –  Requests · Proof Events · Steps",
             Inches(0.2), Inches(5.4), Inches(4.1), Inches(0.3),
             font_size=10, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font_name="Calibri")

add_rect(slide, Inches(4.7), Inches(5.3), Inches(3.9), Inches(0.55), SURFACE_EL, AMBER, 0.6)
add_text_box(slide, "⚡  Redis  –  Circuit Breaker State",
             Inches(4.7), Inches(5.4), Inches(3.9), Inches(0.3),
             font_size=10, bold=True, color=AMBER, align=PP_ALIGN.CENTER, font_name="Calibri")

add_rect(slide, Inches(9.0), Inches(5.3), Inches(3.9), Inches(0.55), SURFACE_EL, PURPLE, 0.6)
add_text_box(slide, "📦  S3 / Blob  –  Backup Storage",
             Inches(9.0), Inches(5.4), Inches(3.9), Inches(0.3),
             font_size=10, bold=True, color=PURPLE, align=PP_ALIGN.CENTER, font_name="Calibri")

# Layer labels (left margin)
layer_labels = [
    (Inches(1.32), "CLIENT"),
    (Inches(2.47), "ORCHESTRATION"),
    (Inches(3.42), "MESSAGING"),
    (Inches(4.32), "WORKERS"),
    (Inches(5.30), "STORAGE"),
]
# (omitted for space — labels already implied by layout)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 · END-TO-END WORKFLOW
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), GREEN)
section_label(slide, "End-to-End Workflow")
heading(slide, "Request lifecycle", y=Inches(0.48))
subheading(slide, "Every deletion request moves through a well-defined state machine with step-level visibility.")

# State machine row
states = [
    ("PENDING",           AMBER,  Inches(0.35)),
    ("RUNNING",           BLUE,   Inches(2.35)),
    ("COMPLETED",         GREEN,  Inches(4.35)),
    ("PARTIAL\nCOMPLETED",BLUE,   Inches(6.35)),
    ("FAILED",            RGBColor(0xd4,0x57,0x57), Inches(4.85)),  # placed below
]
# Draw first two + arrows
for label, color, sx in states[:2]:
    add_rect(slide, sx, Inches(1.72), Inches(1.8), Inches(0.62), SURFACE_EL, color, 0.8)
    add_text_box(slide, label, sx + Inches(0.05), Inches(1.82), Inches(1.7), Inches(0.42),
                 font_size=10, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Calibri")

# Arrow PENDING→RUNNING
arrow_h(slide, Inches(2.15), Inches(2.03), Inches(2.35))
add_text_box(slide, "validate", Inches(2.18), Inches(1.78), Inches(0.5), Inches(0.2),
             font_size=7.5, color=TEXT_MUT, font_name="Calibri")

# From RUNNING branch out to 3 outcomes
# COMPLETED (upper right)
add_rect(slide, Inches(5.1), Inches(1.3), Inches(2.0), Inches(0.62), SURFACE_EL, GREEN, 0.8)
add_text_box(slide, "COMPLETED", Inches(5.1), Inches(1.48), Inches(2.0), Inches(0.28),
             font_size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER, font_name="Calibri")

# PARTIAL_COMPLETED (middle right)
add_rect(slide, Inches(5.1), Inches(2.22), Inches(2.6), Inches(0.62), SURFACE_EL, BLUE, 0.8)
add_text_box(slide, "PARTIAL_COMPLETED", Inches(5.1), Inches(2.4), Inches(2.6), Inches(0.28),
             font_size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Calibri")

# FAILED (lower right)
add_rect(slide, Inches(5.1), Inches(3.14), Inches(2.0), Inches(0.62), SURFACE_EL, RGBColor(0xd4,0x57,0x57), 0.8)
add_text_box(slide, "FAILED", Inches(5.1), Inches(3.32), Inches(2.0), Inches(0.28),
             font_size=10, bold=True, color=RGBColor(0xd4,0x57,0x57), align=PP_ALIGN.CENTER, font_name="Calibri")

# RETRYING (below FAILED)
add_rect(slide, Inches(7.5), Inches(3.14), Inches(2.0), Inches(0.62), SURFACE_EL, PURPLE, 0.8)
add_text_box(slide, "RETRYING", Inches(7.5), Inches(3.32), Inches(2.0), Inches(0.28),
             font_size=10, bold=True, color=PURPLE, align=PP_ALIGN.CENTER, font_name="Calibri")

# DLQ
add_rect(slide, Inches(9.9), Inches(3.14), Inches(1.6), Inches(0.62), SURFACE_EL, RGBColor(0xd4,0x57,0x57), 0.7)
add_text_box(slide, "DLQ", Inches(9.9), Inches(3.32), Inches(1.6), Inches(0.28),
             font_size=10, bold=True, color=RGBColor(0xd4,0x57,0x57), align=PP_ALIGN.CENTER, font_name="Calibri")

# Two detail cards at bottom
card_with_title(slide, "Step-level status per service",
    ["Each service independently reports: PENDING → RUNNING →", "SUCCEEDED / FAILED / RETRYING / SKIPPED_CIRCUIT_OPEN"],
    Inches(0.35), Inches(4.2), Inches(5.8), Inches(1.15), accent_color=BLUE)

card_with_title(slide, "Why PARTIAL_COMPLETED is intentional",
    ["If non-critical services fail but primary cleanup succeeds,", "partial success is recorded. Auditors see exactly which systems",
     "were cleaned — more informative than a binary pass/fail."],
    Inches(6.45), Inches(4.2), Inches(6.55), Inches(1.15), accent_color=GREEN)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 · RELIABILITY DESIGN
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), ORANGE)
section_label(slide, "Reliability Design")
heading(slide, "Built for failure, not against it", y=Inches(0.48))
subheading(slide, "Three layered mechanisms ensure every deletion eventually completes — or is explicitly accounted for.")

rel_cards = [
    ("🔁", "Retry + Dead-Letter Queue", ORANGE,
     ["Failed messages re-queue with exponential backoff",
      "After N attempts, routed to DLQ for manual inspection",
      "DLQ replay endpoint lets ops re-process safely",
      "RabbitMQ TTL + x-dead-letter-exchange wiring"]),
    ("🛡️", "Idempotency Guards", PURPLE,
     ["processed_events table with unique (event_id, service_name)",
      "PostgreSQL 23505 unique violation → service skips silently",
      "Duplicate events logged as DUPLICATE_EVENT_IGNORED",
      "Safe to re-deliver messages at any time"]),
    ("⚡", "Circuit Breaker", BLUE,
     ["Per-service circuit state stored in Redis",
      "CLOSED → OPEN after configurable failure threshold",
      "OPEN steps report SKIPPED_CIRCUIT_OPEN — no wasted retries",
      "Half-open probe auto-resets after cooldown period"]),
]
cx = Inches(0.35)
for icon, title, color, body in rel_cards:
    card_with_title(slide, title, body, cx, Inches(1.75), Inches(4.15), Inches(2.6),
                    icon=icon, accent_color=color, title_color=color)
    cx += Inches(4.35)

# Recovery flow strip
add_rect(slide, Inches(0.35), Inches(4.6), W - Inches(0.7), Inches(0.06), BORDER, None)
add_text_box(slide, "Failure Recovery Flow", Inches(0.35), Inches(4.72), Inches(3), Inches(0.25),
             font_size=9, bold=True, color=TEXT_MUT, font_name="Calibri")

flow2 = [
    ("Worker Fails",       RGBColor(0xd4,0x57,0x57)),
    ("Retry Queue (TTL)",  PURPLE),
    ("Worker Retries",     BLUE),
    ("SUCCEEDED",          GREEN),
]
fx = Inches(0.35)
fy = Inches(5.1)
for label, color in flow2:
    add_rect(slide, fx, fy, Inches(2.2), Inches(0.56), SURFACE_EL, color, 0.7)
    add_text_box(slide, label, fx + Inches(0.1), fy + Inches(0.14), Inches(2.0), Inches(0.3),
                 font_size=10, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Calibri")
    if label != "SUCCEEDED":
        arrow_h(slide, fx + Inches(2.2), fy + Inches(0.28), fx + Inches(2.5))
    fx += Inches(2.55)

# DLQ branch
add_text_box(slide, "max retries → DLQ", Inches(7.45), Inches(5.0), Inches(1.8), Inches(0.2),
             font_size=8.5, color=TEXT_MUT, font_name="Calibri")
add_rect(slide, Inches(7.65), Inches(5.1), Inches(1.6), Inches(0.56), SURFACE_EL, RGBColor(0xd4,0x57,0x57), 0.7)
add_text_box(slide, "DLQ", Inches(7.65), Inches(5.28), Inches(1.6), Inches(0.28),
             font_size=10, bold=True, color=RGBColor(0xd4,0x57,0x57), align=PP_ALIGN.CENTER, font_name="Calibri")
arrow_h(slide, Inches(9.25), Inches(5.38), Inches(9.6))
add_rect(slide, Inches(9.6), Inches(5.1), Inches(1.8), Inches(0.56), SURFACE_EL, AMBER, 0.7)
add_text_box(slide, "Ops Replay API", Inches(9.6), Inches(5.28), Inches(1.8), Inches(0.28),
             font_size=10, bold=True, color=AMBER, align=PP_ALIGN.CENTER, font_name="Calibri")
arrow_h(slide, Inches(11.4), Inches(5.38), Inches(11.75))
add_rect(slide, Inches(11.75), Inches(5.1), Inches(1.4), Inches(0.56), SURFACE_EL, BLUE, 0.7)
add_text_box(slide, "Re-process", Inches(11.75), Inches(5.28), Inches(1.4), Inches(0.28),
             font_size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Calibri")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 · PROOF & AUDITABILITY
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), GREEN)
section_label(slide, "Proof & Auditability")
heading(slide, "Tamper-evident audit chain", y=Inches(0.48))
subheading(slide, "Every deletion event is hash-chained. Modify any record and verification fails — cryptographically.")

pbullets = [
    ("🔗", "Hash-chained proof events",
     "Each event records previous_hash + event_hash (SHA-256). Any tampering breaks the chain."),
    ("🔒", "Advisory lock during append",
     "PostgreSQL pg_advisory_xact_lock ensures sequential per-request chain writes under concurrency."),
    ("📤", "Export & verify",
     "Full proof chain exportable as JSON. One-click 'Verify Chain' traverses all hashes and reports status."),
    ("📋", "Compliance-ready",
     "Deletion evidence satisfies GDPR Art. 17 and CCPA deletion audit requirements out of the box."),
]
by = Inches(1.72)
for icon, title, body in pbullets:
    bullet_card(slide, title, body, Inches(0.35), by, Inches(6.2), icon=icon, accent=GREEN)
    by += Inches(0.97)

# Chain diagram (right side)
chain_items = [
    ("GENESIS",                  GREEN,  "hash: sha256(request_id + 'genesis')\nprev: —"),
    ("DELETION_STEP_SUCCEEDED",  BLUE,   "service: cache-cleanup\nhash: sha256(prev + payload)"),
    ("DELETION_STEP_SUCCEEDED",  PURPLE, "service: search-cleanup\nhash: sha256(prev + payload)"),
    ("DELETION_STEP_SUCCEEDED",  AMBER,  "service: analytics-cleanup\nhash: sha256(prev + payload)"),
]
cy = Inches(1.5)
for label, color, sub in chain_items:
    add_rect(slide, Inches(7.05), cy, Inches(6.0), Inches(0.72), SURFACE_EL, color, 0.7)
    add_text_box(slide, label, Inches(7.15), cy + Inches(0.08), Inches(5.8), Inches(0.26),
                 font_size=10, bold=True, color=color, font_name="Calibri")
    add_text_box(slide, sub, Inches(7.15), cy + Inches(0.34), Inches(5.8), Inches(0.34),
                 font_size=8.5, color=TEXT_MUT, font_name="Calibri")
    cy += Inches(0.86)
    if cy < Inches(1.5) + Inches(0.86) * 3.5:
        # connector
        add_rect(slide, Inches(10.0), cy - Inches(0.14), Inches(0.06), Inches(0.16), GREEN, None)

add_text_box(slide, "↳  … continues for each service …",
             Inches(7.15), cy, Inches(5.8), Inches(0.28),
             font_size=9, color=TEXT_MUT, font_name="Calibri")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 · CONSISTENCY & TRADEOFFS
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), AMBER)
section_label(slide, "Consistency & Tradeoffs")
heading(slide, "Eventual consistency, intentionally", y=Inches(0.48))
subheading(slide, "Not all systems need the same guarantee. We made deliberate, documented tradeoffs per service.")

# Table headers
headers = ["System", "Strategy", "Latency", "Consistency", "Rationale"]
col_widths = [Inches(1.5), Inches(2.0), Inches(1.1), Inches(1.4), Inches(3.5)]
col_x = [Inches(0.35)]
for cw in col_widths[:-1]:
    col_x.append(col_x[-1] + cw + Inches(0.04))

add_rect(slide, Inches(0.35), Inches(1.72), W - Inches(0.7), Inches(0.38), SURFACE_EL, BORDER, 0.5)
for i, h in enumerate(headers):
    add_text_box(slide, h, col_x[i] + Inches(0.06), Inches(1.78), col_widths[i], Inches(0.28),
                 font_size=9, bold=True, color=TEXT_MUT, font_name="Calibri")

rows = [
    ("PostgreSQL",    "Synchronous",       "Low",    "Strong",   "Source of truth; consistent first",     ORANGE),
    ("Cache (Redis)", "Async Event",        "V. Low", "Eventual", "TTL provides natural expiry as backup",  BLUE),
    ("Search Index",  "Async Event",        "Low",    "Eventual", "Search staleness OK for seconds",        BLUE),
    ("Analytics",     "Soft-delete/Delayed","High",   "Eventual", "Batch jobs tolerate delay; immutable logs preferred", PURPLE),
    ("Backup",        "Scheduled",          "High",   "Eventual", "Backup rotation cycle handles deletion", AMBER),
]
ry = Inches(2.1)
for sys, strat, lat, cons, rat, color in rows:
    add_rect(slide, Inches(0.35), ry, W - Inches(0.7), Inches(0.44), SURFACE, color if color else BORDER, 0.4)
    row_data = [sys, strat, lat, cons, rat]
    for i, val in enumerate(row_data):
        bold = (i == 0)
        c = color if i == 1 else (TEXT_PRI if bold else TEXT_SEC)
        add_text_box(slide, val, col_x[i] + Inches(0.06), ry + Inches(0.1), col_widths[i], Inches(0.28),
                     font_size=9.5, bold=bold, color=c, font_name="Calibri")
    ry += Inches(0.48)

# Two insight cards
card_with_title(slide, "Why PARTIAL_COMPLETED is good",
    ["Rather than failing the whole request if analytics is slow, we record",
     "partial success — auditors see exactly which systems were cleaned."],
    Inches(0.35), Inches(4.72), Inches(6.3), Inches(1.1), accent_color=ORANGE)

card_with_title(slide, "CAP Theorem position",
    ["We favour Availability + Partition Tolerance (AP) for the event-driven path.",
     "The proof chain provides Consistency through eventual convergence."],
    Inches(6.95), Inches(4.72), Inches(6.05), Inches(1.1), accent_color=BLUE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 · OBSERVABILITY
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), GREEN)
section_label(slide, "Observability & Operations")
heading(slide, "Production-grade visibility", y=Inches(0.48))
subheading(slide, "Every request, event, and failure is traced, measured, and surfaced — from Jaeger spans to admin dashboards.")

obs_bullets = [
    ("🔭", "OpenTelemetry + Jaeger",
     "Distributed traces show exactly how long each service step took and where failures occurred."),
    ("📊", "Prometheus + Grafana",
     "Request throughput, step duration, retry counts, and circuit breaker state exported as metrics."),
    ("🛠️", "Admin Dashboard",
     "Real-time service health, circuit breaker states, and SLA violation alerts — all in one view."),
    ("⏱️", "SLA Monitoring (SLA-001)",
     "Configurable thresholds per service. Violations recorded with timestamps for rapid triage."),
    ("📡", "Real-time SSE Streaming",
     "Frontend receives live status updates via Server-Sent Events — no polling, no delay."),
]
by = Inches(1.72)
for icon, title, body in obs_bullets:
    bullet_card(slide, title, body, Inches(0.35), by, Inches(6.5), icon=icon, accent=GREEN)
    by += Inches(0.97)

# Stack + admin surfaces
card_with_title(slide, "Observability Stack",
    ["OpenTelemetry · Jaeger · Prometheus · Grafana · SSE",
     "Instrumented at NestJS interceptor level — zero business logic changes."],
    Inches(7.1), Inches(1.72), Inches(5.9), Inches(1.1), accent_color=BLUE)

card_with_title(slide, "Admin Panel Surfaces",
    ["Per-service health status (healthy / degraded / down)",
     "Circuit breaker state per service (CLOSED / OPEN / HALF-OPEN)",
     "SLA violations with timestamps",
     "DLQ message count + replay trigger"],
    Inches(7.1), Inches(3.0), Inches(5.9), Inches(1.5), accent_color=GREEN)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 · DEPLOYMENT
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), BLUE)
section_label(slide, "Deployment")
heading(slide, "Kubernetes-ready, cloud-deployed", y=Inches(0.48))
subheading(slide, "Every service containerised and orchestrated — local via Docker Compose, production-scale on Kubernetes.")

dep_bullets = [
    ("🐳", "Docker Compose (local dev)",
     "Single docker-compose up starts all 10+ services, Postgres, Redis, RabbitMQ, Jaeger, Prometheus."),
    ("☸️", "Kubernetes (production)",
     "Deployment manifests, Services, and ConfigMaps for every component. Horizontal scaling of workers."),
    ("🌐", "Cloud Deployment (GKE)",
     "Live on cloud cluster. All services verified running end-to-end in a production-like environment."),
    ("🔐", "Environment isolation",
     "Secrets managed via .env.example pattern. No credentials committed to source control."),
]
by = Inches(1.72)
for icon, title, body in dep_bullets:
    bullet_card(slide, title, body, Inches(0.35), by, Inches(6.5), icon=icon, accent=BLUE)
    by += Inches(0.97)

# Services tags
add_rect(slide, Inches(7.1), Inches(1.72), Inches(5.9), Inches(2.6), SURFACE_EL, BORDER, 0.5)
add_text_box(slide, "Services in Production", Inches(7.2), Inches(1.82), Inches(5.7), Inches(0.25),
             font_size=9, bold=True, color=TEXT_MUT, font_name="Calibri")
svc_labels = ["Backend API","Frontend","Cache Cleanup","Search Cleanup",
              "Analytics","Backup","Proof Service","RabbitMQ","PostgreSQL","Redis","Jaeger","Prometheus"]
sx, sy = Inches(7.2), Inches(2.12)
for i, svc in enumerate(svc_labels):
    add_rect(slide, sx, sy, Inches(2.6), Inches(0.28), SURFACE, BORDER, 0.4)
    add_text_box(slide, svc, sx + Inches(0.08), sy + Inches(0.04), Inches(2.44), Inches(0.22),
                 font_size=9, color=TEXT_SEC, font_name="Calibri")
    sy += Inches(0.34)
    if i == 5:
        sx = Inches(10.1)
        sy = Inches(2.12)

# Metrics strip
metrics = [("10+", "Microservices", BLUE),("65", "Unit Tests", GREEN),("100%", "Passing", GREEN),("5", "Data Stores", ORANGE)]
mx = Inches(7.1)
my = Inches(4.56)
mw = Inches(1.35)
for val, label, color in metrics:
    add_rect(slide, mx, my, mw, Inches(0.8), SURFACE_EL, color, 0.6)
    add_text_box(slide, val, mx, my + Inches(0.1), mw, Inches(0.36),
                 font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text_box(slide, label, mx, my + Inches(0.46), mw, Inches(0.26),
                 font_size=8, color=TEXT_MUT, align=PP_ALIGN.CENTER, font_name="Calibri")
    mx += Inches(1.5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 · KEY RESULTS
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), ORANGE)
section_label(slide, "Key Results")
heading(slide, "What we built and proved", y=Inches(0.48))

res_cards = [
    ("🏗️", "Complete Distributed System", ORANGE,
     ["10+ microservices coordinating end-to-end deletion", "across 5 data store types — all event-driven."]),
    ("🔒", "Verifiable Proof Chain", BLUE,
     ["Cryptographic hash chain for every deletion.", "Tamper detection works at the single-record level."]),
    ("💪", "Failure-Safe Design", GREEN,
     ["Retry + DLQ + idempotency + circuit breaker means", "no deletion is lost silently — ever."]),
]
cx = Inches(0.35)
for icon, title, color, body in res_cards:
    card_with_title(slide, title, body, cx, Inches(1.7), Inches(4.15), Inches(1.8),
                    icon=icon, accent_color=color, title_color=color)
    cx += Inches(4.35)

# Big metrics
big_metrics = [
    ("65",  "Unit Tests",         BLUE),
    ("15",  "Test Suites",        PURPLE),
    ("5",   "Data Stores",        ORANGE),
    ("3",   "Reliability Layers", GREEN),
    ("10+", "Services Deployed",  AMBER),
]
mx = Inches(0.35)
for val, label, color in big_metrics:
    add_rect(slide, mx, Inches(3.85), Inches(2.4), Inches(1.4), SURFACE_EL, color, 0.7)
    add_text_box(slide, val, mx, Inches(4.0), Inches(2.4), Inches(0.7),
                 font_size=38, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text_box(slide, label, mx, Inches(4.72), Inches(2.4), Inches(0.3),
                 font_size=10, color=TEXT_MUT, align=PP_ALIGN.CENTER, font_name="Calibri")
    mx += Inches(2.55)

# Closing line
add_text_box(slide,
    "EraseGraph demonstrates that compliance-grade deletion propagation is achievable in distributed systems\n"
    "with the right combination of event-driven orchestration, reliability engineering, and cryptographic auditability.",
    Inches(0.35), Inches(5.55), W - Inches(0.7), Inches(0.8),
    font_size=11, color=TEXT_SEC, font_name="Calibri")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 · TEAM CONTRIBUTIONS
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), PURPLE)
section_label(slide, "Team Contributions")
heading(slide, "Who built what", y=Inches(0.48))

# Contribution table headers
add_rect(slide, Inches(0.35), Inches(1.55), W - Inches(0.7), Inches(0.35), SURFACE_EL, BORDER, 0.5)
contrib_headers = ["Member", "Owned Areas", "Key Deliverables", "Impact"]
contrib_col_x = [Inches(0.35), Inches(2.05), Inches(5.55), Inches(9.55)]
contrib_col_w = [Inches(1.6),  Inches(3.4),  Inches(3.9),   Inches(3.45)]
for i, h in enumerate(contrib_headers):
    add_text_box(slide, h, contrib_col_x[i] + Inches(0.08), Inches(1.59), contrib_col_w[i], Inches(0.26),
                 font_size=9, bold=True, color=TEXT_MUT, font_name="Calibri")

contrib_rows = [
    ("Haoyuan Shan", BLUE,
     "GitHub consolidation, K8s / cloud deployment, backend integration, docs coordination",
     "K8s manifests, cloud deployment, integration closure, proof-chain optimisation",
     "Cohesive, demo-ready, submission-ready deliverable"),
    ("Vritika Malhotra", PURPLE,
     "Retry, DLQ, idempotency, circuit breaker, reliability design & documentation",
     "Circuit breaker service, DLQ replay, idempotency guard, reliability spec, unit tests",
     "Failure recovery & robustness of distributed deletion workflows"),
    ("Sakshat Patil", ORANGE,
     "Frontend dashboard, workflow UX, real-time status experience",
     "React pages, SSE streaming UI, bulk upload, admin front-end, demo script",
     "Visualisation quality & demo clarity"),
    ("Asim Mohammed", GREEN,
     "Consistency semantics, backend/admin workflow, status-behaviour refinement",
     "SLA monitor, SSE backend, deletion-request service, state-model design",
     "State-model clarity & operational reliability"),
]

ry = Inches(1.9)
for name, color, areas, deliverables, impact in contrib_rows:
    add_rect(slide, Inches(0.35), ry, W - Inches(0.7), Inches(0.88), SURFACE, color, 0.35)
    row_data = [name, areas, deliverables, impact]
    for i, val in enumerate(row_data):
        bold = (i == 0)
        c = color if bold else TEXT_SEC
        add_text_box(slide, val, contrib_col_x[i] + Inches(0.08), ry + Inches(0.08),
                     contrib_col_w[i] - Inches(0.12), Inches(0.75),
                     font_size=9.5 if i == 0 else 8.5, bold=bold, color=c, font_name="Calibri")
    ry += Inches(0.94)

# RACI legend
add_text_box(slide,
    "R = Responsible   A = Accountable   C = Consulted   I = Informed",
    Inches(0.35), Inches(5.72), W - Inches(0.7), Inches(0.25),
    font_size=8.5, color=TEXT_MUT, align=PP_ALIGN.RIGHT, font_name="Calibri")

# RACI mini-table
raci_headers = ["Area", "Haoyuan", "Vritika", "Sakshat", "Asim"]
raci_rows = [
    ("Backend / API",              "R", "C", "I", "A"),
    ("Frontend",                   "I", "I", "R/A", "C"),
    ("Reliability",                "I", "R/A", "I", "C"),
    ("Proof & Audit Chain",        "R/A", "C", "I", "I"),
    ("Deployment / Cloud",         "R/A", "I", "I", "C"),
    ("Consistency / State Model",  "C", "C", "I", "R/A"),
]
raci_col_x   = [Inches(0.35), Inches(3.35), Inches(5.35), Inches(7.35), Inches(9.35)]
raci_col_w   = [Inches(2.9), Inches(1.9), Inches(1.9), Inches(1.9), Inches(3.0)]
raci_colors  = [BLUE, PURPLE, ORANGE, GREEN]

add_rect(slide, Inches(0.35), Inches(6.0), W - Inches(0.7), Inches(0.28), SURFACE_EL, BORDER, 0.4)
for i, h in enumerate(raci_headers):
    c = TEXT_MUT if i == 0 else raci_colors[i - 1]
    add_text_box(slide, h, raci_col_x[i] + Inches(0.06), Inches(6.02), raci_col_w[i], Inches(0.22),
                 font_size=8, bold=True, color=c, font_name="Calibri")

ry2 = Inches(6.28)
for row in raci_rows:
    add_rect(slide, Inches(0.35), ry2, W - Inches(0.7), Inches(0.18), SURFACE, BORDER, 0.3)
    for i, val in enumerate(row):
        c = TEXT_PRI if i == 0 else raci_colors[i - 1]
        add_text_box(slide, val, raci_col_x[i] + Inches(0.06), ry2 + Inches(0.02),
                     raci_col_w[i], Inches(0.16),
                     font_size=7.5, bold=(i > 0 and "R" in val), color=c, font_name="Calibri")
    ry2 += Inches(0.2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 · Q&A
# ─────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, W, Inches(0.06), ORANGE)

# Large glyph
add_text_box(slide, "🗑️", Inches(5.7), Inches(0.8), Inches(2.0), Inches(1.4),
             font_size=64, align=PP_ALIGN.CENTER, font_name="Segoe UI Emoji")

# Thank you
add_text_box(slide, "Thank you.", Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.0),
             font_size=60, bold=True, color=TEXT_PRI, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text_box(slide, "Questions?", Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.9),
             font_size=52, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font_name="Calibri")

add_text_box(slide,
    "EraseGraph  —  Verifiable deletion propagation in distributed systems",
    Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.4),
    font_size=14, color=TEXT_SEC, align=PP_ALIGN.CENTER, font_name="Calibri")

# Pills row
tags = ["Event-Driven Orchestration", "Proof Chain", "Reliability First", "Compliance Ready", "CMPE 273 · Spring 2026"]
tx = Inches(0.7)
for tag in tags:
    tw = Inches(len(tag) * 0.095 + 0.4)
    add_rect(slide, tx, Inches(4.7), tw, Inches(0.36), SURFACE_EL, BORDER, 0.5)
    add_text_box(slide, tag, tx + Inches(0.1), Inches(4.74), tw - Inches(0.2), Inches(0.26),
                 font_size=9, color=TEXT_SEC, font_name="Calibri")
    tx += tw + Inches(0.18)

# ── Save ────────────────────────────────────────────────────────────────────
out_path = "erasegraph-slides.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({prs.slides.__len__()} slides)")
